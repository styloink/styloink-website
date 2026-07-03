#!/usr/bin/env python3
"""
Extract structured text content from PowerPoint presentations.

This module provides functionality to:
- Extract all text content from PowerPoint shapes
- Preserve paragraph formatting (alignment, bullets, fonts, spacing)
- Handle nested GroupShapes recursively with correct absolute positions
- Sort shapes by visual position on slides
- Filter out slide numbers and non-content placeholders
- Export to JSON with clean, structured data

Classes:
    ParagraphData: Represents a text paragraph with formatting
    ShapeData: Represents a shape with position and text content

Main Functions:
    extract_text_inventory: Extract all text from a presentation
    save_inventory: Save extracted data to JSON

Usage:
    python inventory.py input.pptx output.json
"""

import argparse
import json
import platform
import sys
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape

# Type aliases for cleaner signatures
JsonValue = Union[str, int, float, bool, None]
ParagraphDict = Dict[str, JsonValue]
ShapeDict = Dict[
    str, Union[str, float, bool, List[ParagraphDict], List[str], Dict[str, Any], None]
]
InventoryData = Dict[
    str, Dict[str, "ShapeData"]
]  # Dict of slide_id -> {shape_id -> ShapeData}
InventoryDict = Dict[str, Dict[str, ShapeDict]]  # JSON-serializable inventory


def main():
    """Main entry point for command-line usage."""
    parser = argparse.ArgumentParser(
        description="Extract text inventory from PowerPoint with proper GroupShape support.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python inventory.py presentation.pptx inventory.json
    Extracts text inventory with correct absolute positions for grouped shapes

  python inventory.py presentation.pptx inventory.json --issues-only
    Extracts only text shapes that have overflow or overlap issues

The output JSON includes:
  - All text content organized by slide and shape
  - Correct absolute positions for shapes in groups
  - Visual position and size in inches
  - Paragraph properties and formatting
  - Issue detection: text overflow and shape overlaps
        """,
    )

    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument("output", help="Output JSON file for inventory")
    parser.add_argument(
        "--issues-only",
        action="store_true",
        help="Include only text shapes that have overflow or overlap issues",
    )

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)

    if not input_path.suffix.lower() == ".pptx":
        print("Error: Input must be a PowerPoint file (.pptx)")
        sys.exit(1)

    try:
        print(f"Extracting text inventory from: {args.input}")
        if args.issues_only:
            print(
                "Filtering to include only text shapes with issues (overflow/overlap)"
            )
        inventory = extract_text_inventory(input_path, issues_only=args.issues_only)

        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        save_inventory(inventory, output_path)

        print(f"Output saved to: {args.output}")

        # Report statistics
        total_slides = len(inventory)
        total_shapes = sum(len(shapes) for shapes in inventory.values())
        if args.issues_only:
            if total_shapes > 0:
                print(
                    f"Found {total_shapes} text elements with issues in {total_slides} slides"
                )
            else:
                print("No issues discovered")
        else:
            print(
                f"Found text in {total_slides} slides with {total_shapes} text elements"
            )

    except Exception as e:
        print(f"Error processing presentation: {e}")
        import traceback

        traceback.print_exc()
        sys.exit(1)


@dataclass
class ShapeWithPosition:
    """A shape with its absolute position on the slide."""

    shape: BaseShape
    absolute_left: int  # in EMUs
    absolute_top: int  # in EMUs


class ParagraphData:
    """Data structure for paragraph properties extracted from a PowerPoint paragraph."""

    def __init__(self, paragraph: Any):
        """Initialize from a PowerPoint paragraph object.

        Args:
            paragraph: The PowerPoint paragraph object
        """
        self.text: str = paragraph.text.strip()
        self.bullet: bool = False
        self.level: Optional[int] = None
        self.alignment: Optional[str] = None
        self.space_before: Optional[float] = None
        self.space_after: Optional[float] = None
        self.font_name: Optional[str] = None
        self.font_size: Optional[float] = None
        self.bold: Optional[bool] = None
        self.italic: Optional[bool] = None
        self.underline: Optional[bool] = None
        self.color: Optional[str] = None
        self.theme_color: Optional[str] = None
        self.line_spacing: Optional[float] = None

        # Check for bullet formatting
        if (
            hasattr(paragraph, "_p")
            and paragraph._p is not None
            and paragraph._p.pPr is not None
        ):
            pPr = paragraph._p.pPr
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            if (
                pPr.find(f"{ns}buChar") is not None
                or pPr.find(f"{ns}buAutoNum") is not None
            ):
                self.bullet = True
                if hasattr(paragraph, "level"):
                    self.level = paragraph.level

        # Add alignment if not LEFT (default)
        if hasattr(paragraph, "alignment") and paragraph.alignment is not None:
            alignment_map = {
                PP_ALIGN.CENTER: "CENTER",
                PP_ALIGN.RIGHT: "RIGHT",
                PP_ALIGN.JUSTIFY: "JUSTIFY",
            }
            if paragraph.alignment in alignment_map:
                self.alignment = alignment_map[paragraph.alignment]

        # Add spacing properties if set
        if hasattr(paragraph, "space_before") and paragraph.space_before:
            self.space_before = paragraph.space_before.pt
        if hasattr(paragraph, "space_after") and paragraph.space_after:
            self.space_after = paragraph.space_after.pt

        # Extract font properties from first run
        if paragraph.runs:
            first_run = paragraph.runs[0]
            if hasattr(first_run, "font"):
                font = first_run.font
                if font.name:
                    self.font_name = font.name
                if font.size:
                    self.font_size = font.size.pt
                if font.bold is not None:
                    self.bold = font.bold
                if font.italic is not None:
                    self.italic = font.italic
                if font.underline is not None:
                    self.underline = font.underline

                # Handle color - both RGB and theme colors
                try:
                    # Try RGB color first
                    if font.color.rgb:
                        self.color = str(font.color.rgb)
                except (AttributeError, TypeError):
                    # Fall back to theme color
                    try:
                        if font.color.theme_color:
                            self.theme_color = font.color.theme_color.name
                    except (AttributeError, TypeError):
                        pass

        # Add line spacing if set
        if hasattr(paragraph, "line_spacing") and paragraph.line_spacing is not None:
            if hasattr(paragraph.line_spacing, "pt"):
                self.line_spacing = round(paragraph.line_spacing.pt, 2)
            else:
                # Multiplier - convert to points
                font_size = self.font_size if self.font_size else 12.0
                self.line_spacing = round(paragraph.line_spacing * font_size, 2)

    def to_dict(self) -> ParagraphDict:
        """Convert to dictionary for JSON serialization, excluding None values."""
        result: ParagraphDict = {"text": self.text}

        # Add optional fields only if they have values
        if self.bullet:
            result["bullet"] = self.bullet
        if self.level is not None:
            result["level"] = self.level
        if self.alignment:
            result["alignment"] = self.alignment
        if self.space_before is not None:
            result["space_before"] = self.space_before
        if self.space_after is not None:
            result["space_after"] = self.space_after
        if self.font_name:
            result["font_name"] = self.font_name
        if self.font_size is not None:
            result["font_size"] = self.font_size
        if self.bold is not None:
            result["bold"] = self.bold
        if self.italic is not None:
            result["italic"] = self.italic
        if self.underline is not None:
            result["underline"] = self.underline
        if self.color:
            result["color"] = self.color
        if self.theme_color:
            result["theme_color"] = self.theme_color
        if self.line_spacing is not None:
            result["line_spacing"] = self.line_spacing

        return result


class ShapeData:
    """Data structure for shape properties extracted from a PowerPoint shape."""

    _FONT_PATH_CACHE: Dict[tuple, Optional[str]] = {}
    _CJK_FONT_PATH_CACHE: Optional[str] = None
    _CJK_FONT_PATH_LOOKED_UP = False
    _FALLBACK_FONT_PATH_CACHE: Optional[str] = None
    _FALLBACK_FONT_PATH_LOOKED_UP = False

    # PowerPoint/WPS/Keynote do not share Chromium's exact text metrics. These
    # Small guard bands account for differences between PIL text measurement
    # and PowerPoint's actual rendering engine. Keep them tight — the
    # trailing-leading deduction (see _estimate_frame_overflow) now handles
    # the main over-estimation that larger factors previously compensated for.
    LATIN_WIDTH_FACTOR = 1.02
    # CJK advances are (near-)monospaced full-width boxes that PowerPoint packs
    # tightly. PIL over-measures them; calibrated against decks confirmed to have
    # no visual overflow, raw CJK runs sit a few % wide, so a sub-1.0 factor
    # brings the estimate in line with PowerPoint and avoids false CJK overflow
    # reports while still leaving genuinely overflowing CJK well above threshold.
    CJK_WIDTH_FACTOR = 0.95
    LATIN_LINE_FACTOR = 1.05
    CJK_LINE_FACTOR = 1.05
    # Slack on usable width before deciding a line must wrap. At the ~±2% PIL-vs-
    # PowerPoint metric noise floor, real wraps and fitting text are
    # indistinguishable by width alone, so the tolerance is font-size aware:
    #  - Large fonts (titles): tight tolerance — a title exceeding its one-line
    #    box by ~2% truly wraps in PowerPoint and its wrapped 2nd line is a large,
    #    visible overflow worth flagging.
    #  - Small fonts (labels/footers): looser tolerance — relative metric noise
    #    is higher and a phantom sub-line is visually negligible, so don't wrap.
    # CJK is handled mainly by CJK_WIDTH_FACTOR; it gets a bit more slack here.
    WRAP_TOLERANCE_LARGE_FONT = 1.015   # font_size >= WRAP_LARGE_FONT_PT
    WRAP_TOLERANCE_SMALL_FONT = 1.05
    WRAP_TOLERANCE_CJK = 1.05
    WRAP_LARGE_FONT_PT = 16.0

    @staticmethod
    def emu_to_inches(emu: int) -> float:
        """Convert EMUs (English Metric Units) to inches."""
        return emu / 914400.0

    @staticmethod
    def inches_to_pixels(inches: float, dpi: int = 96) -> int:
        """Convert inches to pixels at given DPI."""
        return int(inches * dpi)

    @staticmethod
    def points_to_pixels(points: float, dpi: int = 96) -> float:
        """Convert points to pixels at given DPI."""
        return points * dpi / 72.0

    @staticmethod
    def _safe_float(value: Any) -> Optional[float]:
        try:
            if value is None:
                return None
            return float(value)
        except (TypeError, ValueError):
            return None

    @staticmethod
    def _emu_attr_to_inches(value: Any) -> Optional[float]:
        numeric = ShapeData._safe_float(value)
        if numeric is None:
            return None
        return numeric / 914400.0

    @staticmethod
    def _length_to_inches(value: Any) -> Optional[float]:
        if value is None:
            return None
        try:
            return float(value) / 914400.0
        except (TypeError, ValueError):
            return None

    @staticmethod
    def _contains_cjk(text: str) -> bool:
        for ch in text:
            code = ord(ch)
            if (
                0x3400 <= code <= 0x9FFF
                or 0xF900 <= code <= 0xFAFF
                or 0x3040 <= code <= 0x30FF
                or 0xAC00 <= code <= 0xD7AF
            ):
                return True
        return False

    @staticmethod
    def _is_cjk_break_char(ch: str) -> bool:
        if ch.isspace():
            return True
        if ShapeData._contains_cjk(ch):
            return True
        # Treat full-width punctuation as a break opportunity. This is a
        # practical approximation of Unicode line breaking for CJK-heavy decks.
        return unicodedata.east_asian_width(ch) in {"W", "F"}

    @staticmethod
    def _tokenize_for_wrap(text: str) -> List[str]:
        tokens: List[str] = []
        current = []
        for ch in text:
            if ch.isspace():
                if current:
                    tokens.append("".join(current))
                    current = []
                tokens.append(" ")
            elif ShapeData._is_cjk_break_char(ch):
                if current:
                    tokens.append("".join(current))
                    current = []
                tokens.append(ch)
            else:
                current.append(ch)
        if current:
            tokens.append("".join(current))
        return tokens

    @staticmethod
    def get_font_path(
        font_name: str, bold: bool = False, italic: bool = False
    ) -> Optional[str]:
        """Get the font file path for a given font name.

        Args:
            font_name: Name of the font (e.g., 'Arial', 'Calibri')
            bold: Prefer the bold face when available (bold text is wider, so
                measuring with the regular face under-estimates width and can
                miss real overflow — e.g. a bold 28pt title that wraps).
            italic: Prefer the italic/oblique face when available.

        Returns:
            Path to the font file, or None if not found
        """
        if not font_name:
            return None
        cache_key = (font_name, bool(bold), bool(italic))
        if cache_key in ShapeData._FONT_PATH_CACHE:
            return ShapeData._FONT_PATH_CACHE[cache_key]

        system = platform.system()

        # Common font file variations to try
        font_variations = [
            font_name,
            font_name.lower(),
            font_name.replace(" ", ""),
            font_name.replace(" ", "-"),
        ]

        # Define font directories and extensions by platform
        if system == "Darwin":  # macOS
            font_dirs = [
                "/System/Library/Fonts/",
                "/System/Library/Fonts/Supplemental/",
                "/Library/Fonts/",
                "~/Library/Fonts/",
            ]
            extensions = [".ttf", ".otf", ".ttc", ".dfont"]
        elif system == "Windows":
            # Windows fonts live under %WINDIR%\Fonts and per-user dirs. Without
            # this branch the lookup fell through to Linux paths, returned None,
            # and the caller used PIL's fixed-size bitmap default (which ignores
            # font size) — causing phantom line wraps and false overflow reports.
            import os as _os

            windir = _os.environ.get("WINDIR", r"C:\Windows")
            localapp = _os.environ.get("LOCALAPPDATA", "")
            font_dirs = [
                _os.path.join(windir, "Fonts"),
                _os.path.join(localapp, "Microsoft", "Windows", "Fonts")
                if localapp
                else "",
            ]
            font_dirs = [d for d in font_dirs if d]
            extensions = [".ttf", ".ttc", ".otf"]
        else:  # Linux
            font_dirs = [
                "/usr/share/fonts/truetype/",
                "/usr/share/fonts/",
                "/usr/local/share/fonts/",
                "~/.fonts/",
            ]
            extensions = [".ttf", ".otf", ".ttc"]

        # Try to find the font file
        from pathlib import Path

        # Tokens that denote a weight/style variant of the SAME family — safe to
        # strip from a candidate stem before comparing. Note: width variants like
        # "narrow"/"condensed" are NOT here — "Arial Narrow" is a different family
        # from "Arial", so stripping them would mis-match. Anything left over that
        # still differs from the target (e.g. "hb", "hebrew", "arabic") means a
        # DIFFERENT typeface and must NOT be treated as a match.
        style_tokens = (
            "regular", "bold", "italic", "oblique", "light", "medium",
            "semibold", "demibold", "black", "heavy", "thin", "book",
            "mt", "ps", "std", "pro",
        )

        target = font_name.lower().replace(" ", "").replace("-", "")

        def _match_rank(stem_lower: str) -> Optional[int]:
            """Return a preference rank (lower = better) if stem matches the
            target family, else None.

            Ranking prefers (1) the requested weight/style (bold/italic) and
            (2) fewer leftover style strips. Using the wrong weight matters:
            bold text is ~8% wider than regular, so measuring a bold title with
            the regular face under-estimates width and misses real wrapping.
            """
            stem = stem_lower.replace(" ", "").replace("-", "").replace("_", "")
            # Detect the candidate's own style from its name.
            has_bold = "bold" in stem or "black" in stem or "heavy" in stem
            has_italic = "italic" in stem or "oblique" in stem

            def _style_penalty() -> int:
                pen = 0
                # Mismatch on the requested attribute is the biggest penalty;
                # an unwanted extra attribute is a smaller one.
                if bool(bold) != has_bold:
                    pen += 4
                if bool(italic) != has_italic:
                    pen += 2
                return pen

            if stem == target:
                # Family name with no style tokens = the regular face.
                base = 0 if (not bold and not italic) else 1
                return base + _style_penalty()

            strips = 0
            work = stem
            changed = True
            while changed and work:
                changed = False
                for tok in style_tokens:
                    if work.endswith(tok) and len(work) > len(tok):
                        work = work[: -len(tok)]
                        strips += 1
                        changed = True
                if work == target:
                    return 1 + strips + _style_penalty()
            return None

        for font_dir in font_dirs:
            font_dir_path = Path(font_dir).expanduser()
            if not font_dir_path.exists():
                continue

            # Exact filename match only when no specific style is requested
            # (an exact base filename is the regular face).
            if not bold and not italic:
                for variant in font_variations:
                    for ext in extensions:
                        font_path = font_dir_path / f"{variant}{ext}"
                        if font_path.exists():
                            ShapeData._FONT_PATH_CACHE[cache_key] = str(font_path)
                            return str(font_path)

            # Strict stem matching: collect candidates whose name (minus
            # weight/style suffixes) EQUALS the target family, and pick the best
            # ranked one (preferring the requested bold/italic face). Recurse one
            # tree because macOS keeps many fonts in Supplemental/.
            best_path: Optional[str] = None
            best_rank = 1 << 30
            try:
                for file_path in font_dir_path.rglob("*"):
                    if not file_path.is_file():
                        continue
                    name_lower = file_path.name.lower()
                    if not any(name_lower.endswith(ext) for ext in extensions):
                        continue
                    stem_lower = name_lower
                    for ext in extensions:
                        if stem_lower.endswith(ext):
                            stem_lower = stem_lower[: -len(ext)]
                            break
                    rank = _match_rank(stem_lower)
                    if rank is not None and rank < best_rank:
                        best_rank = rank
                        best_path = str(file_path)
                        if rank == 0:
                            break  # perfect match — can't do better
            except (OSError, PermissionError):
                continue

            if best_path is not None:
                ShapeData._FONT_PATH_CACHE[cache_key] = best_path
                return best_path

        ShapeData._FONT_PATH_CACHE[cache_key] = None
        return None

    @staticmethod
    def get_cjk_font_path() -> Optional[str]:
        """Return a CJK-capable fallback font when possible."""
        if ShapeData._CJK_FONT_PATH_LOOKED_UP:
            return ShapeData._CJK_FONT_PATH_CACHE

        ShapeData._CJK_FONT_PATH_LOOKED_UP = True
        candidates = [
            "PingFang SC",
            "PingFang",
            "Hiragino Sans GB",
            "STHeiti",
            "Songti SC",
            "Microsoft YaHei",
            "SimSun",
            "Noto Sans CJK SC",
            "Noto Sans CJK",
            "WenQuanYi Zen Hei",
            "Arial Unicode",
        ]
        for name in candidates:
            path = ShapeData.get_font_path(name)
            if path:
                ShapeData._CJK_FONT_PATH_CACHE = path
                return path
        return None

    @staticmethod
    def get_scalable_fallback_font_path() -> Optional[str]:
        """Return a path to ANY scalable (TrueType/OpenType) font.

        Used when the requested font cannot be found. This must NOT fall back to
        ImageFont.load_default(): that bitmap font has a fixed size and ignores
        the requested point size, which makes width measurement meaningless and
        causes phantom line wraps -> false overflow reports (the Windows case).
        """
        if ShapeData._FALLBACK_FONT_PATH_LOOKED_UP:
            return ShapeData._FALLBACK_FONT_PATH_CACHE

        ShapeData._FALLBACK_FONT_PATH_LOOKED_UP = True

        # Try common, metrically-stable sans-serif families across platforms.
        for name in (
            "Arial",
            "Helvetica",
            "DejaVu Sans",
            "Liberation Sans",
            "Verdana",
            "Tahoma",
            "Segoe UI",
            "Calibri",
            "Noto Sans",
        ):
            path = ShapeData.get_font_path(name)
            if path:
                ShapeData._FALLBACK_FONT_PATH_CACHE = path
                return path

        # PIL ships DejaVuSans.ttf — a guaranteed scalable fallback.
        try:
            import PIL

            dejavu = (
                Path(PIL.__file__).parent / "fonts" / "DejaVuSans.ttf"
            )
            if dejavu.exists():
                ShapeData._FALLBACK_FONT_PATH_CACHE = str(dejavu)
                return str(dejavu)
        except Exception:
            pass

        return None

    @staticmethod
    def get_slide_dimensions(slide: Any) -> tuple[Optional[int], Optional[int]]:
        """Get slide dimensions from slide object.

        Args:
            slide: Slide object

        Returns:
            Tuple of (width_emu, height_emu) or (None, None) if not found
        """
        try:
            prs = slide.part.package.presentation_part.presentation
            return prs.slide_width, prs.slide_height
        except (AttributeError, TypeError):
            return None, None

    @staticmethod
    def get_default_font_size(shape: BaseShape, slide_layout: Any) -> Optional[float]:
        """Extract default font size from slide layout for a placeholder shape.

        Args:
            shape: Placeholder shape
            slide_layout: Slide layout containing the placeholder definition

        Returns:
            Default font size in points, or None if not found
        """
        try:
            if not hasattr(shape, "placeholder_format"):
                return None

            shape_type = shape.placeholder_format.type  # type: ignore
            for layout_placeholder in slide_layout.placeholders:
                if layout_placeholder.placeholder_format.type == shape_type:
                    # Find first defRPr element with sz (size) attribute
                    for elem in layout_placeholder.element.iter():
                        if "defRPr" in elem.tag and (sz := elem.get("sz")):
                            return float(sz) / 100.0  # Convert EMUs to points
                    break
        except Exception:
            pass
        return None

    def __init__(
        self,
        shape: BaseShape,
        absolute_left: Optional[int] = None,
        absolute_top: Optional[int] = None,
        slide: Optional[Any] = None,
    ):
        """Initialize from a PowerPoint shape object.

        Args:
            shape: The PowerPoint shape object (should be pre-validated)
            absolute_left: Absolute left position in EMUs (for shapes in groups)
            absolute_top: Absolute top position in EMUs (for shapes in groups)
            slide: Optional slide object to get dimensions and layout information
        """
        self.shape = shape  # Store reference to original shape
        self.shape_id: str = ""  # Will be set after sorting

        # Get slide dimensions from slide object
        self.slide_width_emu, self.slide_height_emu = (
            self.get_slide_dimensions(slide) if slide else (None, None)
        )

        # Get placeholder type if applicable
        self.placeholder_type: Optional[str] = None
        self.default_font_size: Optional[float] = None
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:  # type: ignore
            if shape.placeholder_format and shape.placeholder_format.type:  # type: ignore
                self.placeholder_type = (
                    str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]  # type: ignore
                )

                # Get default font size from layout
                if slide and hasattr(slide, "slide_layout"):
                    self.default_font_size = self.get_default_font_size(
                        shape, slide.slide_layout
                    )

        # Get position information
        # Use absolute positions if provided (for shapes in groups), otherwise use shape's position
        left_emu = (
            absolute_left
            if absolute_left is not None
            else (shape.left if hasattr(shape, "left") else 0)
        )
        top_emu = (
            absolute_top
            if absolute_top is not None
            else (shape.top if hasattr(shape, "top") else 0)
        )

        self.left: float = round(self.emu_to_inches(left_emu), 2)  # type: ignore
        self.top: float = round(self.emu_to_inches(top_emu), 2)  # type: ignore
        self.width: float = round(
            self.emu_to_inches(shape.width if hasattr(shape, "width") else 0),
            2,  # type: ignore
        )
        self.height: float = round(
            self.emu_to_inches(shape.height if hasattr(shape, "height") else 0),
            2,  # type: ignore
        )

        # Store EMU positions for overflow calculations
        self.left_emu = left_emu
        self.top_emu = top_emu
        self.width_emu = shape.width if hasattr(shape, "width") else 0
        self.height_emu = shape.height if hasattr(shape, "height") else 0

        # Calculate overflow status
        self.frame_overflow_bottom: Optional[float] = None
        self.frame_overflow_right: Optional[float] = None
        self.slide_overflow_right: Optional[float] = None
        self.slide_overflow_bottom: Optional[float] = None
        self.estimated_text_left: Optional[float] = None
        self.estimated_text_top: Optional[float] = None
        self.estimated_text_width: Optional[float] = None
        self.estimated_text_height: Optional[float] = None
        self.overlapping_shapes: Dict[
            str, float
        ] = {}  # Dict of shape_id -> overlap area in sq inches
        self.warnings: List[str] = []
        self._estimate_frame_overflow()
        self._calculate_slide_overflow()
        self._detect_bullet_issues()

    @property
    def paragraphs(self) -> List[ParagraphData]:
        """Calculate paragraphs from the shape's text frame."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return []

        paragraphs = []
        for paragraph in self.shape.text_frame.paragraphs:  # type: ignore
            if paragraph.text.strip():
                paragraphs.append(ParagraphData(paragraph))
        return paragraphs

    def _get_default_font_size(self) -> int:
        """Get default font size from theme text styles or use conservative default."""
        try:
            if not (
                hasattr(self.shape, "part") and hasattr(self.shape.part, "slide_layout")
            ):
                return 14

            slide_master = self.shape.part.slide_layout.slide_master  # type: ignore
            if not hasattr(slide_master, "element"):
                return 14

            # Determine theme style based on placeholder type
            style_name = "bodyStyle"  # Default
            if self.placeholder_type and "TITLE" in self.placeholder_type:
                style_name = "titleStyle"

            # Find font size in theme styles
            for child in slide_master.element.iter():
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == style_name:
                    for elem in child.iter():
                        if "sz" in elem.attrib:
                            return int(elem.attrib["sz"]) // 100
        except Exception:
            pass

        return 14  # Conservative default for body text

    def _get_usable_dimensions(self, text_frame) -> Tuple[int, int]:
        """Get usable width and height in pixels after accounting for margins."""
        margins = self._get_text_margins(text_frame)

        usable_width = self.width - margins["left"] - margins["right"]
        usable_height = self.height - margins["top"] - margins["bottom"]

        # Convert to pixels
        return (
            self.inches_to_pixels(usable_width),
            self.inches_to_pixels(usable_height),
        )

    def _get_text_margins(self, text_frame) -> Dict[str, float]:
        """Read textbox insets from OOXML bodyPr, falling back to python-pptx.

        bodyPr stores lIns/rIns/tIns/bIns in EMUs. Reading it directly matters
        because zero-valued insets are common in html2pptx output and older
        truthy checks accidentally treated 0 as "missing".
        """
        margins = {"top": 0.05, "bottom": 0.05, "left": 0.1, "right": 0.1}

        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        body_pr = None
        try:
            body_pr = self.shape.element.find(f".//{ns_a}bodyPr")
        except Exception:
            body_pr = None

        attr_map = {
            "lIns": "left",
            "rIns": "right",
            "tIns": "top",
            "bIns": "bottom",
        }
        ooxml_margin_keys = set()
        if body_pr is not None:
            for attr, key in attr_map.items():
                parsed = self._emu_attr_to_inches(body_pr.get(attr))
                if parsed is not None:
                    margins[key] = parsed
                    ooxml_margin_keys.add(key)

        fallback_props = {
            "margin_top": "top",
            "margin_bottom": "bottom",
            "margin_left": "left",
            "margin_right": "right",
        }
        for prop, key in fallback_props.items():
            if key in ooxml_margin_keys:
                continue
            if hasattr(text_frame, prop):
                parsed = self._length_to_inches(getattr(text_frame, prop))
                if parsed is not None:
                    margins[key] = parsed

        return margins

    def _spacing_points_from_ooxml(
        self, pPr: Any, tag: str, font_size: float
    ) -> Optional[float]:
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        try:
            node = pPr.find(f"{ns_a}{tag}") if pPr is not None else None
            if node is None:
                return None
            pts = node.find(f"{ns_a}spcPts")
            if pts is not None and pts.get("val") is not None:
                return float(pts.get("val")) / 100.0
            pct = node.find(f"{ns_a}spcPct")
            if pct is not None and pct.get("val") is not None:
                return font_size * (float(pct.get("val")) / 100000.0)
        except (TypeError, ValueError):
            return None
        return None

    def _paragraph_indent_inches(self, pPr: Any, is_bullet: bool) -> float:
        """Return the estimated text start offset inside the textbox."""
        mar_l = self._emu_attr_to_inches(pPr.get("marL")) if pPr is not None else None
        indent = self._emu_attr_to_inches(pPr.get("indent")) if pPr is not None else None
        mar_l = mar_l or 0.0
        indent = indent or 0.0

        offset = max(0.0, mar_l) + max(0.0, indent)
        if is_bullet and offset == 0:
            # Common PowerPoint default bullet indent. This is intentionally
            # conservative because missing pPr values still consume line width.
            offset = 0.25
        return offset

    def _wrap_text_line(
        self,
        line: str,
        max_width_px: int,
        draw,
        font,
        width_factor: float,
        wrap_tolerance: float = 1.015,
    ) -> List[str]:
        """Wrap a single line of text to fit within max_width_px."""
        if not line:
            return [""]

        def measure(text: str) -> float:
            try:
                return draw.textlength(text, font=font) * width_factor
            except Exception:
                return len(text) * 7 * width_factor

        # Allow a small slack so text that only marginally exceeds the box is
        # not phantom-wrapped (script-specific; see WRAP_WIDTH_TOLERANCE_*).
        wrap_limit = max_width_px * wrap_tolerance

        # Use textlength for efficient width calculation
        if measure(line) <= wrap_limit:
            return [line]

        wrapped = []
        tokens = self._tokenize_for_wrap(line)
        current_line = ""
        pending_space = ""

        for token in tokens:
            if token == " ":
                if current_line:
                    pending_space = " "
                continue

            test_line = current_line + (pending_space if current_line else "") + token
            if measure(test_line) <= wrap_limit:
                current_line = test_line
                pending_space = ""
            else:
                if current_line:
                    wrapped.append(current_line)
                    current_line = ""
                    pending_space = ""

                if measure(token) <= wrap_limit:
                    current_line = token
                    continue

                # A single token that only marginally exceeds the box (e.g. a
                # short label like "03" or "PMax" measured ~2% over due to PIL
                # vs PowerPoint metric noise) should not be force-split into
                # per-character lines — PowerPoint renders it on one line. Only
                # char-chunk tokens that genuinely cannot fit (long words/URLs).
                if measure(token) <= max_width_px * 1.25:
                    current_line = token
                    continue

                # Last resort for long unbreakable Latin tokens or URLs.
                chunk = ""
                for ch in token:
                    test_chunk = chunk + ch
                    if chunk and measure(test_chunk) > max_width_px:
                        wrapped.append(chunk)
                        chunk = ch
                    else:
                        chunk = test_chunk
                current_line = chunk

        if current_line:
            wrapped.append(current_line)

        return wrapped

    def _get_autofit_mode(self) -> Optional[str]:
        """Return the text autofit mode from bodyPr, or None for fixed-size frames."""
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        try:
            body_pr = self.shape.element.find(f".//{ns_a}bodyPr")
            if body_pr is not None:
                if body_pr.find(f"{ns_a}spAutoFit") is not None:
                    return "spAutoFit"
                if body_pr.find(f"{ns_a}normAutofit") is not None:
                    return "normAutofit"
        except Exception:
            pass
        return None

    def _estimate_frame_overflow(self) -> None:
        """Estimate if text overflows the shape bounds using PIL text measurement."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return

        text_frame = self.shape.text_frame  # type: ignore
        if not text_frame or not text_frame.paragraphs:
            return

        autofit_mode = self._get_autofit_mode()
        if autofit_mode == "spAutoFit":
            return

        margins = self._get_text_margins(text_frame)
        usable_width_px, usable_height_px = self._get_usable_dimensions(text_frame)
        if usable_width_px <= 0 or usable_height_px <= 0:
            return

        # Set up PIL for text measurement
        dummy_img = Image.new("RGB", (1, 1))
        draw = ImageDraw.Draw(dummy_img)

        # Get default font size from placeholder or use conservative estimate
        default_font_size = self._get_default_font_size()

        # word_wrap on (the PowerPoint/html2pptx default) means text reflows to
        # fit the box width instead of spilling right. For such shapes a slightly
        # wide wrapped line is a wrap-estimation artifact, not a real horizontal
        # overflow — only vertical growth matters. word_wrap is None -> inherit
        # default (True); explicit False means single-line that CAN spill right.
        try:
            word_wrap_enabled = text_frame.word_wrap is not False
        except Exception:
            word_wrap_enabled = True

        cursor_y_px = 0.0
        bbox_left = float("inf")
        bbox_top = float("inf")
        bbox_right = float("-inf")
        bbox_bottom = float("-inf")
        last_line_height_px = 0.0
        last_font_size = 0.0
        any_cjk = False
        did_wrap = False
        # A SMALL-FONT single short token (e.g. 7.5pt "PMax") that exceeds its box
        # by a hair is metric noise — measure its width leniently. This must NOT
        # apply to large fonts: a big decorative number like 72pt "03" that does
        # not fit its box genuinely overflows in PowerPoint and must be flagged.
        _ftext = (text_frame.text or "").strip()
        _is_short_token = len(_ftext) <= 6 and not any(c.isspace() for c in _ftext)

        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            if not paragraph.text.strip():
                continue

            para_data = ParagraphData(paragraph)
            pPr = paragraph._p.pPr if hasattr(paragraph, "_p") else None

            font_name = para_data.font_name or "Arial"
            font_size = float(para_data.font_size or default_font_size)
            text = paragraph.text.strip()
            has_cjk = self._contains_cjk(text)
            if has_cjk:
                any_cjk = True
            width_factor = self.CJK_WIDTH_FACTOR if has_cjk else self.LATIN_WIDTH_FACTOR
            if has_cjk:
                wrap_tolerance = self.WRAP_TOLERANCE_CJK
            elif font_size >= self.WRAP_LARGE_FONT_PT:
                wrap_tolerance = self.WRAP_TOLERANCE_LARGE_FONT
            else:
                wrap_tolerance = self.WRAP_TOLERANCE_SMALL_FONT
            line_factor = self.CJK_LINE_FACTOR if has_cjk else self.LATIN_LINE_FACTOR

            is_bold = bool(para_data.bold)
            is_italic = bool(para_data.italic)
            font = None
            font_path = self.get_cjk_font_path() if has_cjk else None
            if font_path is None:
                font_path = self.get_font_path(font_name, bold=is_bold, italic=is_italic)
            if font_path is None:
                # Use a scalable TTF fallback, never the fixed-size bitmap
                # default (which ignores font_size and triggers phantom wraps).
                font_path = self.get_scalable_fallback_font_path()
            # Load the measuring font at the SAME 96-DPI scale used for box
            # dimensions (inches_to_pixels) and line height (points_to_pixels).
            # Loading it at the raw point value treats points as pixels (~72 dpi),
            # under-measuring text width by ~25% — which makes titles that wrap in
            # PowerPoint look like they fit on one line, so real overflow is missed.
            font_px = max(1, int(round(self.points_to_pixels(font_size))))
            if font_path:
                try:
                    font = ImageFont.truetype(font_path, size=font_px)
                except Exception:
                    font = ImageFont.load_default()
            else:
                font = ImageFont.load_default()

            para_indent_in = self._paragraph_indent_inches(pPr, para_data.bullet)
            para_width_px = max(1, usable_width_px - self.inches_to_pixels(para_indent_in))

            # Wrap all lines in this paragraph
            all_wrapped_lines = []
            explicit_lines = text.split("\n")
            for line in explicit_lines:
                wrapped = self._wrap_text_line(
                    line, para_width_px, draw, font, width_factor, wrap_tolerance
                )
                if len(wrapped) > 1:
                    did_wrap = True
                all_wrapped_lines.extend(wrapped)

            if all_wrapped_lines:
                # Calculate line height
                ooxml_line_spacing = self._spacing_points_from_ooxml(
                    pPr, "lnSpc", font_size
                )
                if ooxml_line_spacing is not None:
                    line_height_px = self.points_to_pixels(ooxml_line_spacing)
                elif para_data.line_spacing:
                    line_height_px = self.points_to_pixels(para_data.line_spacing)
                else:
                    line_height_px = self.points_to_pixels(font_size) * line_factor

                # For single-line text within this paragraph, line-spacing
                # only determines inter-line distance (irrelevant here). Cap
                # effective height to font ascent+descent (≈ 1.3× font_size).
                font_size_px = self.points_to_pixels(font_size)
                if len(all_wrapped_lines) == 1:
                    max_single_line_px = font_size_px * 1.3
                    effective_line_height_px = min(line_height_px, max_single_line_px)
                else:
                    effective_line_height_px = line_height_px

                last_line_height_px = effective_line_height_px
                last_font_size = font_size

                # Add space_before (except first paragraph)
                space_before = self._spacing_points_from_ooxml(
                    pPr, "spcBef", font_size
                )
                if space_before is None:
                    space_before = para_data.space_before
                if para_idx > 0 and space_before:
                    cursor_y_px += self.points_to_pixels(space_before)

                x_in = self.left + margins["left"] + para_indent_in
                for line_idx, line in enumerate(all_wrapped_lines):
                    try:
                        line_width_px = draw.textlength(line, font=font) * width_factor
                    except Exception:
                        line_width_px = len(line) * 7 * width_factor
                    y_in = self.top + margins["top"] + cursor_y_px / 96.0
                    line_right = x_in + line_width_px / 96.0
                    # Last line uses effective height; others use full
                    # line-spacing to preserve inter-line distance.
                    lh = (
                        effective_line_height_px
                        if line_idx == len(all_wrapped_lines) - 1
                        else line_height_px
                    )
                    line_bottom = y_in + lh / 96.0

                    bbox_left = min(bbox_left, x_in)
                    bbox_top = min(bbox_top, y_in)
                    bbox_right = max(bbox_right, line_right)
                    bbox_bottom = max(bbox_bottom, line_bottom)
                    cursor_y_px += line_height_px

                # Add space_after
                space_after = self._spacing_points_from_ooxml(
                    pPr, "spcAft", font_size
                )
                if space_after is None:
                    space_after = para_data.space_after
                if space_after:
                    cursor_y_px += self.points_to_pixels(space_after)

        if bbox_left == float("inf"):
            return

        # Deduct trailing leading: PowerPoint does not render line-spacing
        # below the last line — only ascent+descent occupy space there.
        if last_line_height_px > 0 and last_font_size > 0:
            font_rendered_height_px = self.points_to_pixels(last_font_size) * 1.2
            trailing_leading_px = max(0.0, last_line_height_px - font_rendered_height_px)
            if trailing_leading_px > 0:
                bbox_bottom -= trailing_leading_px / 96.0

        self.estimated_text_left = round(bbox_left, 2)
        self.estimated_text_top = round(bbox_top, 2)
        self.estimated_text_width = round(max(0.0, bbox_right - bbox_left), 2)
        self.estimated_text_height = round(max(0.0, bbox_bottom - bbox_top), 2)

        frame_right = self.left + self.width - margins["right"]
        frame_bottom = self.top + self.height - margins["bottom"]

        # CJK text measurement is less precise — use wider tolerance to avoid
        # false-positive overflow reports that waste agent iterations.
        if autofit_mode == "normAutofit":
            thresh_right = 0.06
            thresh_bottom = 0.16
        elif any_cjk:
            thresh_right = 0.05
            thresh_bottom = 0.16
        else:
            thresh_right = 0.03
            thresh_bottom = 0.08

        # Horizontal tolerance combines a small absolute floor with a fraction
        # of the box width. PIL-vs-PowerPoint metric error scales with the text
        # run, so a fixed 0.03in floor over-flags tiny decorative labels that
        # exceed the box by a hair while under-tolerating wide titles. Using
        # max(floor, 5% of width) tolerates ~one-glyph noise but still flags a
        # title whose text clearly spills (and would wrap) past the frame.
        usable_width = max(0.0, self.width - margins["left"] - margins["right"])
        # Proportional slack mirroring the wrap tolerance (font-size aware) so a
        # single-line title (or large decorative number) that truly spills is
        # flagged, while small-font / CJK metric noise is not. The short-token
        # leniency applies only to SMALL fonts — a large "03" that overflows is
        # a real defect.
        small_font = last_font_size < self.WRAP_LARGE_FONT_PT
        if any_cjk or (_is_short_token and small_font):
            right_slack = 0.05
        elif not small_font:
            right_slack = 0.015
        else:
            right_slack = 0.05
        thresh_right = max(thresh_right, usable_width * right_slack)

        # When word-wrap is on and the text actually wrapped, horizontal fit is
        # already guaranteed by reflow — a wide wrapped line is wrap-point noise,
        # not a spill. Report only vertical growth in that case.
        suppress_right = word_wrap_enabled and did_wrap

        overflow_right = round(bbox_right - frame_right, 2)
        if overflow_right > thresh_right and not suppress_right:
            self.frame_overflow_right = overflow_right

        overflow_bottom = round(bbox_bottom - frame_bottom, 2)
        if overflow_bottom > thresh_bottom:
            self.frame_overflow_bottom = overflow_bottom

    def overlap_rect(self) -> Tuple[float, float, float, float]:
        """Return estimated rendered text bbox, falling back to shape bbox."""
        if (
            self.estimated_text_left is not None
            and self.estimated_text_top is not None
            and self.estimated_text_width is not None
            and self.estimated_text_height is not None
        ):
            return (
                self.estimated_text_left,
                self.estimated_text_top,
                self.estimated_text_width,
                self.estimated_text_height,
            )
        return (self.left, self.top, self.width, self.height)

    def _calculate_slide_overflow(self) -> None:
        """Calculate if shape overflows the slide boundaries."""
        if self.slide_width_emu is None or self.slide_height_emu is None:
            return

        # Check right overflow (ignore negligible overflows <= 0.01")
        right_edge_emu = self.left_emu + self.width_emu
        if right_edge_emu > self.slide_width_emu:
            overflow_emu = right_edge_emu - self.slide_width_emu
            overflow_inches = round(self.emu_to_inches(overflow_emu), 2)
            if overflow_inches > 0.01:  # Only report significant overflows
                self.slide_overflow_right = overflow_inches

        # Check bottom overflow (ignore negligible overflows <= 0.01")
        bottom_edge_emu = self.top_emu + self.height_emu
        if bottom_edge_emu > self.slide_height_emu:
            overflow_emu = bottom_edge_emu - self.slide_height_emu
            overflow_inches = round(self.emu_to_inches(overflow_emu), 2)
            if overflow_inches > 0.01:  # Only report significant overflows
                self.slide_overflow_bottom = overflow_inches

    def _detect_bullet_issues(self) -> None:
        """Detect bullet point formatting issues in paragraphs."""
        if not self.shape or not hasattr(self.shape, "text_frame"):
            return

        text_frame = self.shape.text_frame  # type: ignore
        if not text_frame or not text_frame.paragraphs:
            return

        # Common bullet symbols that indicate manual bullets
        bullet_symbols = ["•", "●", "○"]

        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            # Check for manual bullet symbols
            if text and any(text.startswith(symbol + " ") for symbol in bullet_symbols):
                self.warnings.append(
                    "manual_bullet_symbol: use proper bullet formatting"
                )
                break

    @property
    def has_any_issues(self) -> bool:
        """Check if shape has any issues (overflow, overlap, or warnings)."""
        return (
            self.frame_overflow_bottom is not None
            or self.frame_overflow_right is not None
            or self.slide_overflow_right is not None
            or self.slide_overflow_bottom is not None
            or len(self.overlapping_shapes) > 0
            or len(self.warnings) > 0
        )

    def to_dict(self) -> ShapeDict:
        """Convert to dictionary for JSON serialization."""
        result: ShapeDict = {
            "left": self.left,
            "top": self.top,
            "width": self.width,
            "height": self.height,
        }

        # Add optional fields if present
        if self.placeholder_type:
            result["placeholder_type"] = self.placeholder_type

        if self.default_font_size:
            result["default_font_size"] = self.default_font_size

        # Add overflow information only if there is overflow
        overflow_data = {}

        # Add frame overflow if present
        if self.frame_overflow_bottom is not None:
            overflow_data["frame"] = {"overflow_bottom": self.frame_overflow_bottom}
        if self.frame_overflow_right is not None:
            overflow_data.setdefault("frame", {})[
                "overflow_right"
            ] = self.frame_overflow_right

        # Add slide overflow if present
        slide_overflow = {}
        if self.slide_overflow_right is not None:
            slide_overflow["overflow_right"] = self.slide_overflow_right
        if self.slide_overflow_bottom is not None:
            slide_overflow["overflow_bottom"] = self.slide_overflow_bottom
        if slide_overflow:
            overflow_data["slide"] = slide_overflow

        # Only add overflow field if there is overflow
        if overflow_data:
            result["overflow"] = overflow_data

        # Add overlap field if there are overlapping shapes
        if self.overlapping_shapes:
            result["overlap"] = {"overlapping_shapes": self.overlapping_shapes}

        # Add warnings field if there are warnings
        if self.warnings:
            result["warnings"] = self.warnings

        if self.estimated_text_left is not None:
            result["estimated_text_bbox"] = {
                "left": self.estimated_text_left,
                "top": self.estimated_text_top,
                "width": self.estimated_text_width,
                "height": self.estimated_text_height,
            }

        # Add paragraphs after placeholder_type
        result["paragraphs"] = [para.to_dict() for para in self.paragraphs]

        return result


def is_valid_shape(shape: BaseShape) -> bool:
    """Check if a shape contains meaningful text content."""
    # Must have a text frame with content
    if not hasattr(shape, "text_frame") or not shape.text_frame:  # type: ignore
        return False

    text = shape.text_frame.text.strip()  # type: ignore
    if not text:
        return False

    # Skip slide numbers and numeric footers
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:  # type: ignore
        if shape.placeholder_format and shape.placeholder_format.type:  # type: ignore
            placeholder_type = (
                str(shape.placeholder_format.type).split(".")[-1].split(" ")[0]  # type: ignore
            )
            if placeholder_type == "SLIDE_NUMBER":
                return False
            if placeholder_type == "FOOTER" and text.isdigit():
                return False

    return True


def collect_shapes_with_absolute_positions(
    shape: BaseShape, parent_left: int = 0, parent_top: int = 0
) -> List[ShapeWithPosition]:
    """Recursively collect all shapes with valid text, calculating absolute positions.

    For shapes within groups, their positions are relative to the group.
    This function calculates the absolute position on the slide by accumulating
    parent group offsets.

    Args:
        shape: The shape to process
        parent_left: Accumulated left offset from parent groups (in EMUs)
        parent_top: Accumulated top offset from parent groups (in EMUs)

    Returns:
        List of ShapeWithPosition objects with absolute positions
    """
    if hasattr(shape, "shapes"):  # GroupShape
        result = []
        # Get this group's position
        group_left = shape.left if hasattr(shape, "left") else 0
        group_top = shape.top if hasattr(shape, "top") else 0

        # Calculate absolute position for this group
        abs_group_left = parent_left + group_left
        abs_group_top = parent_top + group_top

        # Process children with accumulated offsets
        for child in shape.shapes:  # type: ignore
            result.extend(
                collect_shapes_with_absolute_positions(
                    child, abs_group_left, abs_group_top
                )
            )
        return result

    # Regular shape - check if it has valid text
    if is_valid_shape(shape):
        # Calculate absolute position
        shape_left = shape.left if hasattr(shape, "left") else 0
        shape_top = shape.top if hasattr(shape, "top") else 0

        return [
            ShapeWithPosition(
                shape=shape,
                absolute_left=parent_left + shape_left,
                absolute_top=parent_top + shape_top,
            )
        ]

    return []


def sort_shapes_by_position(shapes: List[ShapeData]) -> List[ShapeData]:
    """Sort shapes by visual position (top-to-bottom, left-to-right).

    Shapes within 0.5 inches vertically are considered on the same row.
    """
    if not shapes:
        return shapes

    # Sort by top position first
    shapes = sorted(shapes, key=lambda s: (s.top, s.left))

    # Group shapes by row (within 0.5 inches vertically)
    result = []
    row = [shapes[0]]
    row_top = shapes[0].top

    for shape in shapes[1:]:
        if abs(shape.top - row_top) <= 0.5:
            row.append(shape)
        else:
            # Sort current row by left position and add to result
            result.extend(sorted(row, key=lambda s: s.left))
            row = [shape]
            row_top = shape.top

    # Don't forget the last row
    result.extend(sorted(row, key=lambda s: s.left))
    return result


def calculate_overlap(
    rect1: Tuple[float, float, float, float],
    rect2: Tuple[float, float, float, float],
    tolerance: float = 0.05,
) -> Tuple[bool, float]:
    """Calculate if and how much two rectangles overlap.

    Args:
        rect1: (left, top, width, height) of first rectangle in inches
        rect2: (left, top, width, height) of second rectangle in inches
        tolerance: Minimum overlap in inches to consider as overlapping (default: 0.05")

    Returns:
        Tuple of (overlaps, overlap_area) where:
        - overlaps: True if rectangles overlap by more than tolerance
        - overlap_area: Area of overlap in square inches
    """
    left1, top1, w1, h1 = rect1
    left2, top2, w2, h2 = rect2

    # Calculate overlap dimensions
    overlap_width = min(left1 + w1, left2 + w2) - max(left1, left2)
    overlap_height = min(top1 + h1, top2 + h2) - max(top1, top2)

    # Check if there's meaningful overlap (more than tolerance)
    if overlap_width > tolerance and overlap_height > tolerance:
        # Calculate overlap area in square inches
        overlap_area = overlap_width * overlap_height
        return True, round(overlap_area, 2)

    return False, 0


def detect_overlaps(shapes: List[ShapeData]) -> None:
    """Detect overlapping shapes and update their overlapping_shapes dictionaries.

    This function requires each ShapeData to have its shape_id already set.
    It modifies the shapes in-place, adding shape IDs with overlap areas in square inches.

    Args:
        shapes: List of ShapeData objects with shape_id attributes set
    """
    n = len(shapes)

    # Compare each pair of shapes
    for i in range(n):
        for j in range(i + 1, n):
            shape1 = shapes[i]
            shape2 = shapes[j]

            # Ensure shape IDs are set
            assert shape1.shape_id, f"Shape at index {i} has no shape_id"
            assert shape2.shape_id, f"Shape at index {j} has no shape_id"

            rect1 = shape1.overlap_rect()
            rect2 = shape2.overlap_rect()

            overlaps, overlap_area = calculate_overlap(rect1, rect2)

            if overlaps and overlap_area >= 0.15:
                shape1.overlapping_shapes[shape2.shape_id] = overlap_area
                shape2.overlapping_shapes[shape1.shape_id] = overlap_area


def extract_text_inventory(
    pptx_path: Path, prs: Optional[Any] = None, issues_only: bool = False
) -> InventoryData:
    """Extract text content from all slides in a PowerPoint presentation.

    Args:
        pptx_path: Path to the PowerPoint file
        prs: Optional Presentation object to use. If not provided, will load from pptx_path.
        issues_only: If True, only include shapes that have overflow or overlap issues

    Returns a nested dictionary: {slide-N: {shape-N: ShapeData}}
    Shapes are sorted by visual position (top-to-bottom, left-to-right).
    The ShapeData objects contain the full shape information and can be
    converted to dictionaries for JSON serialization using to_dict().
    """
    if prs is None:
        prs = Presentation(str(pptx_path))
    inventory: InventoryData = {}

    for slide_idx, slide in enumerate(prs.slides):
        # Collect all valid shapes from this slide with absolute positions
        shapes_with_positions = []
        for shape in slide.shapes:  # type: ignore
            shapes_with_positions.extend(collect_shapes_with_absolute_positions(shape))

        if not shapes_with_positions:
            continue

        # Convert to ShapeData with absolute positions and slide reference
        shape_data_list = [
            ShapeData(
                swp.shape,
                swp.absolute_left,
                swp.absolute_top,
                slide,
            )
            for swp in shapes_with_positions
        ]

        # Sort by visual position and assign stable IDs in one step
        sorted_shapes = sort_shapes_by_position(shape_data_list)
        for idx, shape_data in enumerate(sorted_shapes):
            shape_data.shape_id = f"shape-{idx}"

        # Detect overlaps using the stable shape IDs
        if len(sorted_shapes) > 1:
            detect_overlaps(sorted_shapes)

        # Filter for issues only if requested (after overlap detection)
        if issues_only:
            sorted_shapes = [sd for sd in sorted_shapes if sd.has_any_issues]

        if not sorted_shapes:
            continue

        # Create slide inventory using the stable shape IDs
        inventory[f"slide-{slide_idx}"] = {
            shape_data.shape_id: shape_data for shape_data in sorted_shapes
        }

    return inventory


def get_inventory_as_dict(pptx_path: Path, issues_only: bool = False) -> InventoryDict:
    """Extract text inventory and return as JSON-serializable dictionaries.

    This is a convenience wrapper around extract_text_inventory that returns
    dictionaries instead of ShapeData objects, useful for testing and direct
    JSON serialization.

    Args:
        pptx_path: Path to the PowerPoint file
        issues_only: If True, only include shapes that have overflow or overlap issues

    Returns:
        Nested dictionary with all data serialized for JSON
    """
    inventory = extract_text_inventory(pptx_path, issues_only=issues_only)

    # Convert ShapeData objects to dictionaries
    dict_inventory: InventoryDict = {}
    for slide_key, shapes in inventory.items():
        dict_inventory[slide_key] = {
            shape_key: shape_data.to_dict() for shape_key, shape_data in shapes.items()
        }

    return dict_inventory


def save_inventory(inventory: InventoryData, output_path: Path) -> None:
    """Save inventory to JSON file with proper formatting.

    Converts ShapeData objects to dictionaries for JSON serialization.
    """
    # Convert ShapeData objects to dictionaries
    json_inventory: InventoryDict = {}
    for slide_key, shapes in inventory.items():
        json_inventory[slide_key] = {
            shape_key: shape_data.to_dict() for shape_key, shape_data in shapes.items()
        }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(json_inventory, f, indent=2, ensure_ascii=False)


if __name__ == "__main__":
    main()
