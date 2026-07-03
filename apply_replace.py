import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_replacements(input_pptx, replacements_json, output_pptx):
    with open(replacements_json, 'r', encoding='utf-8') as f:
        replacements = json.load(f)
    
    prs = Presentation(input_pptx)
    
    for slide_key, shapes_data in replacements.items():
        slide_idx = int(slide_key.split('-')[1])
        slide = prs.slides[slide_idx]
        
        # We need to map shape index carefully
        # In inventory.py, shapes are ordered by top-to-bottom, left-to-right
        shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
        
        for shape_key, shape_data in shapes_data.items():
            shape_idx = int(shape_key.split('-')[1])
            if shape_idx >= len(shapes):
                continue
            
            shape = shapes[shape_idx]
            if not shape.has_text_frame:
                continue
            
            tf = shape.text_frame
            tf.clear()
            
            for para_data in shape_data.get('paragraphs', []):
                p = tf.add_paragraph()
                p.text = para_data.get('text', '')
                
                if 'alignment' in para_data:
                    p.alignment = getattr(PP_ALIGN, para_data['alignment'])
                
                if para_data.get('bullet'):
                    p.level = para_data.get('level', 0)
                
                run = p.runs[0] if p.runs else p.add_run()
                if 'font_size' in para_data:
                    run.font.size = Pt(para_data['font_size'])
                if para_data.get('bold'):
                    run.font.bold = True
                if 'color' in para_data:
                    # Convert hex to RGB
                    hex_color = para_data['color']
                    run.font.color.rgb = RGBColor.from_string(hex_color)
                # Note: theme_color is complex to set via python-pptx without deep XML, 
                # skipping for now or using hardcoded RGB if known.
    
    prs.save(output_pptx)

if __name__ == "__main__":
    import sys
    apply_replacements(sys.argv[1], sys.argv[2], sys.argv[3])
