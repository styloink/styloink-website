import sys
from pptx import Presentation

def extract_text(pptx_path, output_path):
    prs = Presentation(pptx_path)
    with open(output_path, 'w', encoding='utf-8') as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")
            f.write("\n")

if __name__ == "__main__":
    extract_text(sys.argv[1], sys.argv[2])
